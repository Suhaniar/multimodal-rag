from reportlab.pdfgen import canvas
from docx import Document
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw

import os

os.makedirs("test_data", exist_ok=True)

# ---------------- PDF ----------------
def create_pdf():
    c = canvas.Canvas("test_data/financial_report.pdf")
    c.drawString(100, 800, "Financial Report 2025")
    c.drawString(100, 780, "Q1 Revenue: $2.5M")
    c.drawString(100, 760, "Q2 Revenue: $3.1M")
    c.drawString(100, 740, "Q3 Revenue: $3.8M")
    c.drawString(100, 720, "Q4 Revenue: $4.2M")
    c.drawString(100, 700, "Revenue increased 68% YoY")
    c.save()

# ---------------- CSV ----------------
def create_csv():
    df = pd.DataFrame({
        "Quarter": ["Q1", "Q2", "Q3", "Q4"],
        "Revenue": [2.5, 3.1, 3.8, 4.2]
    })
    df.to_csv("test_data/sales.csv", index=False)

# ---------------- DOCX ----------------
def create_docx():
    doc = Document()
    doc.add_heading("HR Policy", 0)
    doc.add_paragraph("Employees are entitled to 24 days of leave per year.")
    doc.add_paragraph("Work from home is allowed 2 days per week.")
    doc.save("test_data/policy.docx")

# ---------------- PPTX ----------------
def create_ppt():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Results"
    slide.placeholders[1].text = "Revenue increased steadily from Q1 to Q4"
    prs.save("test_data/presentation.pptx")

# ---------------- IMAGE (OCR test) ----------------
def create_image():
    img = Image.new("RGB", (400, 200), color="white")
    d = ImageDraw.Draw(img)
    d.text((50, 80), "Invoice Total: $1234", fill="black")
    img.save("test_data/invoice.png")

if __name__ == "__main__":
    create_pdf()
    create_csv()
    create_docx()
    create_ppt()
    create_image()
    print("Test files created in /test_data")