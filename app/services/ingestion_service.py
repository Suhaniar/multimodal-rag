import uuid
import csv
import io
import shutil
import tempfile
from pathlib import Path
from typing import Any, List, Optional, Tuple, Dict
from datetime import datetime
import pymupdf as fitz
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from app.core.config import (
    CHUNK_SIZE, CHUNK_OVERLAP, OCR_MIN_NATIVE_TEXT_CHARS,
    OCR_MAX_WORKERS, PREVIEW_CHARS, SUPPORTED_EXTENSIONS, OCR_ZOOM,
    OLLAMA_LLM
)
from app.helpers.text_utils import (
    clean_text, rows_to_text, document_key, deduplicate_documents,
    detect_headings, assign_sections_to_lines, extract_headings_from_pdf,
    infer_heading_level, build_section_tree
)
from app.services.ocr_service import run_ocr, enrich_image, understand_chart
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed
import ollama

logger = logging.getLogger("multimodal_rag")

# ---------- Helper to create a Document object ----------
def make_document(
    content: str,
    source: str,
    modality: str,
    kind: str,
    document_id: str | None = None,
    parent_document_id: str | None = None,
    section: str = "",
    section_path: str = "",
    subsection: str = "",
    paragraph_num: int | None = None,
    **metadata: Any
) -> Document | None:
    cleaned = clean_text(content)
    if not cleaned:
        return None

    document_id = document_id or str(uuid.uuid4())
    page = metadata.get("page")
    if page is None:
        page = metadata.get("slide", 0)

    return Document(
        page_content=cleaned,
        metadata={
            "document_id": document_id,
            "parent_document_id": parent_document_id or document_id,
            "chunk_id": None,  # set later during chunking
            "source": source,
            "file_type": Path(source).suffix.lower().replace(".", ""),
            "modality": modality,
            "kind": kind,
            "page": page,
            "parent_page": page,
            "slide": metadata.get("slide"),
            "section": section,
            "section_path": section_path,   # new: full hierarchical path
            "subsection": subsection,
            "paragraph_num": paragraph_num,
            "created_at": datetime.utcnow().isoformat(),
            **metadata,
        },
    )

# ---------- PDF extraction functions (with section hierarchy) ----------
def extract_pdf_text_documents(pdf_path: Path, name: str, document_id: str) -> list[Document]:
    docs = []
    pdf = fitz.open(pdf_path)

    # 1. Collect all headings across pages
    all_headings = []  # (page_num, line_index, heading_text)
    for page_num in range(len(pdf)):
        page = pdf[page_num]
        text = page.get_text("text")
        headings = detect_headings(text)
        for idx, head in headings:
            all_headings.append((page_num + 1, idx, head))

    # 2. Build global section paths (per page)
    # We'll build a dict: page_num -> section_path_string
    section_per_page = {}
    if all_headings:
        # Sort by page, then line index
        all_headings.sort(key=lambda x: (x[0], x[1]))
        path = []
        current_section = ""
        for page_num, idx, heading in all_headings:
            level = infer_heading_level(heading)
            while len(path) >= level:
                path.pop()
            path.append(heading)
            current_section = " > ".join(path)
            section_per_page[page_num] = current_section
        # Propagate to pages without headings: carry previous section
        prev = ""
        for p in range(1, len(pdf) + 1):
            if p in section_per_page:
                prev = section_per_page[p]
            else:
                section_per_page[p] = prev
    else:
        # No headings: all pages get empty section path
        for p in range(1, len(pdf) + 1):
            section_per_page[p] = ""

    # 3. Extract text per page and assign section_path
    for page_num, page in enumerate(pdf, start=1):
        text = page.get_text("text")
        if not text:
            continue
        doc = make_document(
            text,
            source=name,
            modality="text",
            kind="pdf_page_text",
            document_id=document_id,
            page=page_num,
            section_path=section_per_page.get(page_num, ""),
        )
        if doc:
            docs.append(doc)
    pdf.close()
    return docs

def extract_pdf_ocr_documents(pdf_path: Path, name: str, document_id: str, skip_pages: set[int] | None = None) -> list[Document]:
    docs = []
    skip_pages = skip_pages or set()
    try:
        pdf = fitz.open(pdf_path)
    except Exception as e:
        logger.error(f"Failed to open PDF for OCR {name}: {e}", exc_info=True)
        raise ValueError(f"Corrupt or unreadable PDF: {name}")

    # Build section paths for all pages (reuse logic from text extraction)
    all_headings = []
    for p in range(len(pdf)):
        page = pdf[p]
        text = page.get_text("text")
        headings = detect_headings(text)
        for idx, head in headings:
            all_headings.append((p + 1, idx, head))
    section_per_page = {}
    if all_headings:
        all_headings.sort(key=lambda x: (x[0], x[1]))
        path = []
        current_section = ""
        for page_num, idx, heading in all_headings:
            level = infer_heading_level(heading)
            while len(path) >= level:
                path.pop()
            path.append(heading)
            current_section = " > ".join(path)
            section_per_page[page_num] = current_section
        prev = ""
        for p in range(1, len(pdf) + 1):
            if p in section_per_page:
                prev = section_per_page[p]
            else:
                section_per_page[p] = prev
    else:
        for p in range(1, len(pdf) + 1):
            section_per_page[p] = ""

    for page_num, page in enumerate(pdf, start=1):
        if page_num in skip_pages:
            continue
        matrix = fitz.Matrix(OCR_ZOOM, OCR_ZOOM)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
        ocr_text = run_ocr(image)
        doc = make_document(
            ocr_text,
            source=name,
            modality="ocr",
            kind="pdf_page_ocr",
            document_id=document_id,
            page=page_num,
            section_path=section_per_page.get(page_num, ""),
        )
        if doc:
            docs.append(doc)
    pdf.close()
    return docs

def process_single_image_for_doc(
    page_index: int,
    image_index: int,
    width: int,
    height: int,
    extension: str,
    pil_image: Image.Image,
    name: str,
    document_id: str,
    section_path: str = ""
) -> Optional[Document]:
    try:
        ocr_text = run_ocr(pil_image)
        enrichment = enrich_image(pil_image)
        chart_description = understand_chart(pil_image)  # Phase 2: chart understanding

        parts = []
        if enrichment.get("caption"):
            parts.append(f"Caption: {enrichment['caption']}")
        if enrichment.get("summary"):
            parts.append(f"Scene: {enrichment['summary']}")
        if enrichment.get("objects"):
            parts.append(f"Objects: {', '.join(enrichment['objects'])}")
        if ocr_text:
            parts.append(f"Text: {ocr_text}")
        if chart_description:
            parts.append(f"Chart description: {chart_description}")

        combined_content = "\n".join(parts) if parts else f"Image on page {page_index}, image {image_index}. No descriptive content extracted."

        metadata = {
            "page": page_index,
            "image_index": image_index,
            "width": width,
            "height": height,
            "extension": extension,
            "caption": enrichment.get("caption", ""),
            "objects": enrichment.get("objects", []),
            "summary": enrichment.get("summary", ""),
            "ocr_text": ocr_text,
            "chart_description": chart_description,
            "section_path": section_path,
        }
        return make_document(
            combined_content,
            source=name,
            modality="image",
            kind="pdf_embedded_image",
            document_id=document_id,
            page=page_index,
            section_path=section_path,
            **metadata
        )
    except Exception as e:
        logger.error(f"Error processing image on page {page_index}: {e}", exc_info=True)
        return None

def extract_pdf_image_documents(pdf_path: Path, name: str, document_id: str) -> list[Document]:
    docs = []
    try:
        pdf = fitz.open(pdf_path)
    except Exception as e:
        logger.error(f"Failed to open PDF for image extraction {name}: {e}", exc_info=True)
        raise ValueError(f"Corrupt or unreadable PDF: {name}")

    # Build section paths (same as before)
    all_headings = []
    for p in range(len(pdf)):
        page = pdf[p]
        text = page.get_text("text")
        headings = detect_headings(text)
        for idx, head in headings:
            all_headings.append((p + 1, idx, head))
    section_per_page = {}
    if all_headings:
        all_headings.sort(key=lambda x: (x[0], x[1]))
        path = []
        current_section = ""
        for page_num, idx, heading in all_headings:
            level = infer_heading_level(heading)
            while len(path) >= level:
                path.pop()
            path.append(heading)
            current_section = " > ".join(path)
            section_per_page[page_num] = current_section
        prev = ""
        for p in range(1, len(pdf) + 1):
            if p in section_per_page:
                prev = section_per_page[p]
            else:
                section_per_page[p] = prev
    else:
        for p in range(1, len(pdf) + 1):
            section_per_page[p] = ""

    image_tasks = []
    for page_index, page in enumerate(pdf, start=1):
        section_path = section_per_page.get(page_index, "")
        for image_index, image_data in enumerate(page.get_images(full=True), start=1):
            xref = image_data[0]
            image_info = pdf.extract_image(xref)
            width = image_info.get("width")
            height = image_info.get("height")
            extension = image_info.get("ext", "image")
            try:
                pil_image = Image.open(io.BytesIO(image_info["image"])).convert("RGB")
                image_tasks.append((page_index, image_index, width, height, extension, pil_image, section_path))
            except Exception as e:
                logger.warning(f"Could not open image on page {page_index}, image {image_index}: {e}")
                fallback_content = f"Image on page {page_index}, image {image_index}. Format: {extension}. Size: {width}x{height}. No text detected."
                doc = make_document(
                    fallback_content,
                    source=name,
                    modality="image",
                    kind="pdf_embedded_image",
                    document_id=document_id,
                    page=page_index,
                    image_index=image_index,
                    width=width,
                    height=height,
                    extension=extension,
                    section_path=section_path,
                )
                if doc:
                    docs.append(doc)
    pdf.close()

    if not image_tasks:
        return docs

    logger.info(f"Processing {len(image_tasks)} images (OCR + enrichment + chart) with {OCR_MAX_WORKERS} workers")
    with ThreadPoolExecutor(max_workers=OCR_MAX_WORKERS) as executor:
        future_to_task = {}
        for task in image_tasks:
            page_idx, img_idx, w, h, ext, pil_img, sec_path = task
            future = executor.submit(
                process_single_image_for_doc,
                page_idx, img_idx, w, h, ext, pil_img, name, document_id, sec_path
            )
            future_to_task[future] = task
        for future in as_completed(future_to_task):
            doc = future.result()
            if doc:
                docs.append(doc)
    logger.info(f"Completed image processing for {len(image_tasks)} images")
    return docs

# ---------- Table summarization (Phase 2) ----------
def summarize_table(table_text: str) -> str:
    """Generate a one‑sentence summary of the table content."""
    if not table_text.strip():
        return ""
    prompt = f"Summarize the following table in one sentence, focusing on key numbers and trends:\n{table_text}"
    try:
        response = ollama.chat(
            model=OLLAMA_LLM,
            messages=[{"role": "user", "content": prompt}]
        )
        summary = response["message"]["content"].strip()
        return summary
    except Exception as e:
        logger.error(f"Table summarization failed: {e}")
        return ""

# ---------- UPDATED: Clean table extraction with section paths ----------
def extract_pdf_table_documents(pdf_path: Path, name: str, document_id: str) -> list[Document]:
    docs = []

    # ----- Build section paths per page (using fitz) -----
    try:
        pdf_fitz = fitz.open(pdf_path)
        all_headings = []
        for p in range(len(pdf_fitz)):
            page = pdf_fitz[p]
            text = page.get_text("text")
            headings = detect_headings(text)
            for idx, head in headings:
                all_headings.append((p + 1, idx, head))

        section_per_page = {}
        if all_headings:
            all_headings.sort(key=lambda x: (x[0], x[1]))
            path = []
            current_section = ""
            for page_num, idx, heading in all_headings:
                level = infer_heading_level(heading)
                while len(path) >= level:
                    path.pop()
                path.append(heading)
                current_section = " > ".join(path)
                section_per_page[page_num] = current_section
            # Propagate to pages without headings
            prev = ""
            for p in range(1, len(pdf_fitz) + 1):
                if p in section_per_page:
                    prev = section_per_page[p]
                else:
                    section_per_page[p] = prev
        else:
            for p in range(1, len(pdf_fitz) + 1):
                section_per_page[p] = ""
        pdf_fitz.close()
    except Exception as e:
        logger.warning(f"Could not build section paths for {name}: {e}")
        section_per_page = {}

    # ----- Extract tables with pdfplumber -----
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                section_path = section_per_page.get(page_index, "")
                tables = page.extract_tables()
                for table_index, table in enumerate(tables, start=1):
                    if not table or len(table) == 0:
                        continue
                    table_text = rows_to_text(table)

                    # Main table document
                    table_doc = make_document(
                        table_text,
                        source=name,
                        modality="table",
                        kind="pdf_table",
                        document_id=document_id,
                        page=page_index,
                        table_index=table_index,
                        section_path=section_path,
                    )
                    if table_doc:
                        docs.append(table_doc)

                    # Summary document (Phase 2)
                    summary = summarize_table(table_text)
                    if summary:
                        summary_doc = make_document(
                            summary,
                            source=name,
                            modality="text",
                            kind="table_summary",
                            document_id=document_id,
                            page=page_index,
                            table_index=table_index,
                            section_path=section_path,
                            summary_for_table=table_index,
                        )
                        if summary_doc:
                            docs.append(summary_doc)
    except Exception as e:
        logger.error(f"Failed to extract tables from PDF {name}: {e}", exc_info=True)

    return docs

def parse_pdf_file(pdf_path: Path, name: str, document_id: str) -> list[Document]:
    try:
        text_docs = extract_pdf_text_documents(pdf_path, name, document_id)
        skip_pages = {doc.metadata["page"] for doc in text_docs if len(doc.page_content) >= OCR_MIN_NATIVE_TEXT_CHARS}
        ocr_docs = extract_pdf_ocr_documents(pdf_path, name, document_id, skip_pages)
        image_docs = extract_pdf_image_documents(pdf_path, name, document_id)
        table_docs = extract_pdf_table_documents(pdf_path, name, document_id)
        docs = text_docs + ocr_docs + image_docs + table_docs
        if not docs:
            raise ValueError("No content could be extracted from the PDF.")
        return docs
    except ValueError as ve:
        raise ve
    except Exception as e:
        logger.error(f"Unexpected error parsing PDF {name}: {e}", exc_info=True)
        raise ValueError(f"Failed to parse PDF: {name}. Error: {str(e)}")

# ---------- Other file types (updated with section_path for consistency) ----------
def extract_image_file_documents(file_path: Path, name: str, document_id: str) -> list[Document]:
    with Image.open(file_path) as image:
        image = image.convert("RGB")
        width, height = image.size
        image_format = image.format or file_path.suffix.lstrip(".")
        ocr_text = run_ocr(image)
        enrichment = enrich_image(image)
        chart_desc = understand_chart(image)  # Phase 2

        parts = []
        if enrichment.get("caption"):
            parts.append(f"Caption: {enrichment['caption']}")
        if enrichment.get("summary"):
            parts.append(f"Scene: {enrichment['summary']}")
        if enrichment.get("objects"):
            parts.append(f"Objects: {', '.join(enrichment['objects'])}")
        if ocr_text:
            parts.append(f"Text: {ocr_text}")
        if chart_desc:
            parts.append(f"Chart description: {chart_desc}")

        content = "\n".join(parts) if parts else f"Image '{name}' ({image_format}, {width}x{height}). No descriptive content extracted."
        metadata = {
            "width": width,
            "height": height,
            "format": image_format,
            "page": 1,
            "caption": enrichment.get("caption", ""),
            "objects": enrichment.get("objects", []),
            "summary": enrichment.get("summary", ""),
            "ocr_text": ocr_text,
            "chart_description": chart_desc,
            "section_path": "",
        }
        doc = make_document(
            content,
            source=name,
            modality="image",
            kind="uploaded_image",
            document_id=document_id,
            page=1,
            section_path="",
            **metadata
        )
        return [doc] if doc else []

def extract_docx_documents(file_path: Path, name: str, document_id: str) -> list[Document]:
    docs = []
    docx = DocxDocument(file_path)

    # Extract headings with paragraph index
    headings = []
    for i, para in enumerate(docx.paragraphs):
        if para.style.name and para.style.name.startswith('Heading'):
            headings.append((i, para.text.strip()))

    # Build section paths for paragraphs (simple: last heading)
    section_map = {}
    if headings:
        current = ""
        for i in range(len(docx.paragraphs)):
            # if this index is a heading, update current
            for h_idx, h_text in headings:
                if h_idx == i:
                    current = h_text
                    break
            section_map[i] = current
    else:
        for i in range(len(docx.paragraphs)):
            section_map[i] = ""

    # Process paragraphs
    for para_idx, para in enumerate(docx.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        section = section_map.get(para_idx, "")
        doc = make_document(
            text,
            source=name,
            modality="text",
            kind="docx_paragraph",
            document_id=document_id,
            page=1,
            section=section,
            section_path=section,
            paragraph_num=para_idx,
        )
        if doc:
            docs.append(doc)

    # Process tables (no heading propagation, use empty section_path)
    for table_index, table in enumerate(docx.tables, start=1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_text = rows_to_text(rows)
        table_doc = make_document(
            table_text,
            source=name,
            modality="table",
            kind="docx_table",
            document_id=document_id,
            page=1,
            table_index=table_index,
            section_path="",
        )
        if table_doc:
            docs.append(table_doc)
        # Summarize table
        summary = summarize_table(table_text)
        if summary:
            summary_doc = make_document(
                summary,
                source=name,
                modality="text",
                kind="table_summary",
                document_id=document_id,
                page=1,
                table_index=table_index,
                section_path="",
                summary_for_table=table_index,
            )
            if summary_doc:
                docs.append(summary_doc)
    return docs

def extract_pptx_documents(file_path: Path, name: str, document_id: str) -> list[Document]:
    docs = []
    presentation = Presentation(file_path)
    slide_images = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        # Get slide title (first text shape) as section
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip()
                break
        text_parts = []
        image_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
            if shape.shape_type == 13:  # picture
                image_count += 1
                try:
                    pil_image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    slide_images.append((slide_index, shape.shape_id, pil_image, title))
                except Exception as e:
                    logger.warning(f"Could not open image on slide {slide_index}: {e}")
        # Document for slide text
        text_doc = make_document(
            "\n".join(text_parts),
            source=name,
            modality="text",
            kind="pptx_slide_text",
            document_id=document_id,
            slide=slide_index,
            page=slide_index,
            section=title,
            section_path=title,
        )
        if text_doc:
            docs.append(text_doc)
        # Metadata for images count
        if image_count:
            image_doc = make_document(
                f"Slide {slide_index} contains {image_count} image(s).",
                source=name,
                modality="image",
                kind="pptx_slide_image_metadata",
                document_id=document_id,
                slide=slide_index,
                page=slide_index,
                image_count=image_count,
                section=title,
                section_path=title,
            )
            if image_doc:
                docs.append(image_doc)

    # Process images in parallel
    if slide_images:
        logger.info(f"Processing {len(slide_images)} PPTX images (OCR + enrichment + chart)")
        with ThreadPoolExecutor(max_workers=OCR_MAX_WORKERS) as executor:
            future_to_img = {}
            for slide_idx, shape_id, pil_img, sec in slide_images:
                future = executor.submit(lambda img: (run_ocr(img), enrich_image(img), understand_chart(img)), pil_img)
                future_to_img[future] = (slide_idx, shape_id, sec)
            for future in as_completed(future_to_img):
                slide_idx, shape_id, sec = future_to_img[future]
                try:
                    ocr_text, enrichment, chart_desc = future.result()
                    parts = []
                    if enrichment.get("caption"):
                        parts.append(f"Caption: {enrichment['caption']}")
                    if enrichment.get("summary"):
                        parts.append(f"Scene: {enrichment['summary']}")
                    if enrichment.get("objects"):
                        parts.append(f"Objects: {', '.join(enrichment['objects'])}")
                    if ocr_text:
                        parts.append(f"Text: {ocr_text}")
                    if chart_desc:
                        parts.append(f"Chart description: {chart_desc}")
                    content = "\n".join(parts) if parts else f"Image on slide {slide_idx}. No descriptive content."
                    metadata = {
                        "slide": slide_idx,
                        "shape_id": shape_id,
                        "caption": enrichment.get("caption", ""),
                        "objects": enrichment.get("objects", []),
                        "summary": enrichment.get("summary", ""),
                        "ocr_text": ocr_text,
                        "chart_description": chart_desc,
                        "page": slide_idx,
                        "section": sec,
                        "section_path": sec,
                    }
                    doc = make_document(
                        content,
                        source=name,
                        modality="image",
                        kind="pptx_slide_image",
                        document_id=document_id,
                        slide=slide_idx,
                        page=slide_idx,
                        section=sec,
                        section_path=sec,
                        **metadata
                    )
                    if doc:
                        docs.append(doc)
                except Exception as e:
                    logger.error(f"Failed to process PPTX image on slide {slide_idx}: {e}")
    return docs

def extract_plain_text_documents(file_path: Path, name: str, document_id: str) -> list[Document]:
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    # For plain text, we could attempt to detect headings (e.g., markdown) but keep simple.
    doc = make_document(
        text,
        source=name,
        modality="text",
        kind="plain_text",
        document_id=document_id,
        page=1,
        section_path="",
    )
    return [doc] if doc else []

def extract_csv_documents(file_path: Path, name: str, document_id: str) -> list[Document]:
    with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as file:
        rows = list(csv.reader(file))
    table_text = rows_to_text(rows)
    doc = make_document(
        table_text,
        source=name,
        modality="table",
        kind="csv_table",
        document_id=document_id,
        page=1,
        section_path="",
    )
    docs = [doc] if doc else []
    # Summarize CSV table
    summary = summarize_table(table_text)
    if summary:
        summary_doc = make_document(
            summary,
            source=name,
            modality="text",
            kind="table_summary",
            document_id=document_id,
            page=1,
            section_path="",
            summary_for_table=1,
        )
        if summary_doc:
            docs.append(summary_doc)
    return docs

# ---------- Chunking (parent‑child hierarchy) ----------
def chunk_documents(documents: list[Document]) -> list[Document]:
    documents = deduplicate_documents(documents)
    splittable = [doc for doc in documents if doc.metadata["modality"] in {"text", "ocr", "table"}]
    keep = [doc for doc in documents if doc.metadata["modality"] == "image"]

    all_chunks = []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )

    parent_docs = []
    for doc in splittable:
        # Create a parent document with the full content
        parent_meta = doc.metadata.copy()
        parent_meta["chunk_id"] = None
        parent_meta["kind"] = "document"  # mark as parent
        parent_doc = Document(
            page_content=doc.page_content,
            metadata=parent_meta
        )
        parent_docs.append(parent_doc)

        # Now split into chunks
        chunks = splitter.split_documents([doc])
        for idx, chunk in enumerate(chunks):
            # Inherit metadata from parent
            for key, value in parent_meta.items():
                if key not in chunk.metadata:
                    chunk.metadata[key] = value
            chunk.metadata["chunk_id"] = f"{doc.metadata['document_id']}_chunk_{idx}"
            chunk.metadata["chunk_index"] = idx
            chunk.metadata["parent_document_id"] = doc.metadata["document_id"]
            all_chunks.append(chunk)

    all_chunks = deduplicate_documents(all_chunks)
    return parent_docs + all_chunks + keep

def parse_upload(file_path: Path, name: str, document_id: str) -> list[Document]:
    if file_path.stat().st_size == 0:
        raise ValueError(f"File '{name}' is empty.")

    extension = file_path.suffix.lower()
    try:
        if extension == ".pdf":
            docs = parse_pdf_file(file_path, name, document_id)
        elif extension in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
            docs = extract_image_file_documents(file_path, name, document_id)
        elif extension == ".docx":
            docs = extract_docx_documents(file_path, name, document_id)
        elif extension == ".pptx":
            docs = extract_pptx_documents(file_path, name, document_id)
        elif extension in {".txt", ".md"}:
            docs = extract_plain_text_documents(file_path, name, document_id)
        elif extension == ".csv":
            docs = extract_csv_documents(file_path, name, document_id)
        else:
            raise ValueError(f"Unsupported file type '{extension}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
    except Exception as e:
        logger.error(f"Error parsing {name}: {e}", exc_info=True)
        raise ValueError(f"Failed to process file '{name}': {str(e)}")

    if not docs:
        raise ValueError(f"No indexable content could be extracted from '{name}'.")

    return chunk_documents(docs)